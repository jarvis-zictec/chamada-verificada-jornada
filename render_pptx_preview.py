from pathlib import Path
import io, shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

ROOT=Path(__file__).resolve().parent
OUT=ROOT/'artifacts'/'chamada_verificada_collateral'
PPTX=OUT/'zictec-chamada-verificada-apresentacao.pptx'
PREVIEW_DIR=OUT/'previews'
PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
CACHE=Path('/root/.hermes/document_cache')
prs=Presentation(str(PPTX))
W,H=1600,900
SX=W/prs.slide_width; SY=H/prs.slide_height
try:
    fonts={
        'big': ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf', 34),
        'med': ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf', 23),
        'small': ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf', 17),
        'tiny': ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf', 12),
    }
except Exception:
    fonts={k:None for k in ['big','med','small','tiny']}

def rgb_tuple(rgb, default):
    try:
        if rgb is None: return default
        return (rgb[0],rgb[1],rgb[2])
    except Exception:
        return default

def shape_fill(sh, default=None):
    try:
        if sh.fill and sh.fill.type:
            return rgb_tuple(sh.fill.fore_color.rgb, default)
    except Exception:
        pass
    return default

def text_color(sh, default=(23,23,41)):
    try:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                c = r.font.color.rgb
                if c: return rgb_tuple(c, default)
    except Exception:
        pass
    return default

def draw_wrapped(d, text, xy, max_w, max_h, font, fill):
    x,y=xy
    if not text: return
    avg=8 if font==fonts['tiny'] else 10 if font==fonts['small'] else 13
    max_chars=max(8,int(max_w/avg))
    words=text.replace('\x0b',' ').split(); lines=[]; cur=''
    for word in words:
        if len(cur)+len(word)+1>max_chars:
            if cur: lines.append(cur)
            cur=word
        else:
            cur=(cur+' '+word).strip()
    if cur: lines.append(cur)
    line_h=15 if font==fonts['tiny'] else 22 if font==fonts['small'] else 28
    for ln in lines[:max(1,int(max_h/line_h))]:
        d.text((x,y), ln, fill=fill, font=font)
        y += line_h

previews=[]
for idx, slide in enumerate(prs.slides,1):
    im=Image.new('RGB',(W,H),(251,248,246)); d=ImageDraw.Draw(im)
    # Pass 1: filled shapes / backgrounds
    for sh in slide.shapes:
        x=int(sh.left*SX); y=int(sh.top*SY); w=int(sh.width*SX); h=int(sh.height*SY)
        if w<=1 or h<=1: continue
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            fill=shape_fill(sh)
            if fill is not None:
                # rounded approx for cards, rectangle for bg
                if x <= 2 and y <= 2 and w >= W*0.9 and h >= H*0.9:
                    d.rectangle([0,0,W,H], fill=fill)
                else:
                    try:
                        d.rounded_rectangle([x,y,x+w,y+h], radius=18, fill=fill, outline=(235,228,223), width=1)
                    except Exception:
                        d.rectangle([x,y,x+w,y+h], fill=fill)
    # Pass 2: pictures
    for sh in slide.shapes:
        x=int(sh.left*SX); y=int(sh.top*SY); w=int(sh.width*SX); h=int(sh.height*SY)
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                pic=Image.open(io.BytesIO(sh.image.blob)).convert('RGBA')
                pic.thumbnail((w,h))
                px=x+(w-pic.width)//2; py=y+(h-pic.height)//2
                im.paste(pic,(px,py),pic if pic.mode=='RGBA' else None)
            except Exception:
                pass
    # Pass 3: text
    for sh in slide.shapes:
        if not hasattr(sh,'text') or not sh.text.strip(): continue
        x=int(sh.left*SX); y=int(sh.top*SY); w=int(sh.width*SX); h=int(sh.height*SY)
        text=sh.text.strip()
        if h > 80 and w > 650 and len(text) < 120:
            font=fonts['big']
        elif len(text)<50:
            font=fonts['med']
        elif len(text)<190:
            font=fonts['small']
        else:
            font=fonts['tiny']
        draw_wrapped(d, text, (x+8,y+8), w-16, h-12, font, text_color(sh))
    out=PREVIEW_DIR/f'slide-{idx:02d}.png'
    im.save(out); previews.append(out)
thumbs=[Image.open(p).resize((320,180)) for p in previews]
sheet=Image.new('RGB',(960,720),(240,240,240)); dr=ImageDraw.Draw(sheet)
for i,t in enumerate(thumbs):
    x=(i%3)*320; y=(i//3)*180
    sheet.paste(t,(x,y)); dr.text((x+6,y+6),str(i+1),fill=(255,90,31),font=fonts['med'])
contact=PREVIEW_DIR/'pptx-contact-sheet.png'
sheet.save(contact)
shutil.copy2(contact,CACHE/contact.name)
print(contact)
