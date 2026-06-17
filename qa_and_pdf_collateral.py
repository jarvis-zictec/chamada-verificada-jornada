from pathlib import Path
import re, zipfile, shutil
from xml.etree import ElementTree as ET

from docx import Document
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image as RLImage
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics

ROOT = Path(__file__).resolve().parent
OUT = ROOT / 'artifacts' / 'chamada_verificada_collateral'
CACHE = Path('/root/.hermes/document_cache')
PPTX = OUT / 'zictec-chamada-verificada-apresentacao.pptx'
DOCX = OUT / 'zictec-chamada-verificada-datasheet-whitepaper.docx'
MD = OUT / 'zictec-chamada-verificada-datasheet-whitepaper.md'
PDF = OUT / 'zictec-chamada-verificada-datasheet-whitepaper.pdf'
PREVIEW_DIR = OUT / 'previews'
PREVIEW_DIR.mkdir(parents=True, exist_ok=True)

# PDF generation from canonical MD, with explicit tables and corporate styling
md = MD.read_text(encoding='utf-8')
styles = getSampleStyleSheet()
styles.add(ParagraphStyle(name='CoverTitle', parent=styles['Title'], fontSize=24, leading=30, textColor=colors.HexColor('#1B1648'), alignment=1, spaceAfter=16))
styles.add(ParagraphStyle(name='Kicker', parent=styles['Normal'], fontSize=9, leading=12, textColor=colors.HexColor('#FF5A1F'), alignment=1, spaceAfter=10))
styles.add(ParagraphStyle(name='BodyJust', parent=styles['BodyText'], fontSize=9.4, leading=13.5, alignment=4, spaceAfter=5))
styles.add(ParagraphStyle(name='Small', parent=styles['BodyText'], fontSize=8, leading=10, textColor=colors.HexColor('#66657A')))
styles.add(ParagraphStyle(name='H1Z', parent=styles['Heading1'], fontSize=15, leading=19, textColor=colors.HexColor('#1B1648'), spaceBefore=10, spaceAfter=6))
styles.add(ParagraphStyle(name='H2Z', parent=styles['Heading2'], fontSize=12, leading=16, textColor=colors.HexColor('#1B1648'), spaceBefore=8, spaceAfter=4))
styles.add(ParagraphStyle(name='BulletZ', parent=styles['BodyText'], fontSize=9.2, leading=12.5, leftIndent=12, bulletIndent=4, spaceAfter=3))

story=[]
logo = ROOT / 'assets' / 'logo_orange.png'
if logo.exists():
    story.append(RLImage(str(logo), width=4.0*cm, height=1.1*cm))
    story[-1].hAlign='CENTER'
story.append(Paragraph('DATASHEET / WHITEPAPER TÉCNICO-COMERCIAL', styles['Kicker']))
story.append(Paragraph('Chamada Verificada, STIR/SHAKEN e Origem Verificada', styles['CoverTitle']))
story.append(Paragraph('Guia de ofertas, arquitetura comercial, matriz de decisão e premissas para implantação assistida.', styles['BodyJust']))
story.append(Spacer(1, 0.5*cm))
meta=Table([['Fonte','Página Chamada Verificada ZICTEC'],['Uso','Pré-venda, reunião executiva e alinhamento técnico'],['Preparado por','Jarvis / ZICTEC']], colWidths=[3.2*cm, 11.8*cm])
meta.setStyle(TableStyle([('BACKGROUND',(0,0),(0,-1),colors.HexColor('#1B1648')),('TEXTCOLOR',(0,0),(0,-1),colors.white),('BACKGROUND',(1,0),(1,-1),colors.HexColor('#FFF1EA')),('GRID',(0,0),(-1,-1),0.5,colors.HexColor('#E7E0DC')),('FONTNAME',(0,0),(-1,-1),'Helvetica'),('FONTSIZE',(0,0),(-1,-1),9),('VALIGN',(0,0),(-1,-1),'MIDDLE'),('LEFTPADDING',(0,0),(-1,-1),8),('RIGHTPADDING',(0,0),(-1,-1),8),('TOPPADDING',(0,0),(-1,-1),7),('BOTTOMPADDING',(0,0),(-1,-1),7)]))
story.append(meta)
story.append(PageBreak())

def add_table(data, widths):
    t=Table(data, colWidths=widths, repeatRows=1)
    t.setStyle(TableStyle([
        ('BACKGROUND',(0,0),(-1,0),colors.HexColor('#1B1648')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
        ('GRID',(0,0),(-1,-1),0.35,colors.HexColor('#E7E0DC')),('VALIGN',(0,0),(-1,-1),'TOP'),('FONTSIZE',(0,0),(-1,-1),7.2),
        ('LEFTPADDING',(0,0),(-1,-1),4),('RIGHTPADDING',(0,0),(-1,-1),4),('TOPPADDING',(0,0),(-1,-1),4),('BOTTOMPADDING',(0,0),(-1,-1),4),
        ('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white, colors.HexColor('#FFF8F5')])
    ]))
    story.append(t); story.append(Spacer(1,0.25*cm))

lines=md.splitlines()
i=0
while i < len(lines):
    line=lines[i].strip()
    if not line or line.startswith('**Fonte') or line.startswith('**Data') or line.startswith('**Uso'):
        i+=1; continue
    if line.startswith('# '):
        # represented on cover
        i+=1; continue
    if line.startswith('## '):
        story.append(Paragraph(line[3:], styles['H1Z'])); i+=1; continue
    if line.startswith('### '):
        story.append(Paragraph(line[4:], styles['H2Z'])); i+=1; continue
    if line.startswith('|'):
        block=[]
        while i < len(lines) and lines[i].strip().startswith('|'):
            block.append(lines[i].strip()); i+=1
        table_lines=[b for b in block if not re.match(r'^\|\s*-+', b)]
        rows=[]
        for b in table_lines:
            cells=[c.strip().replace('**','') for c in b.strip('|').split('|')]
            rows.append([Paragraph(c, styles['Small']) for c in cells])
        if rows:
            col_count=len(rows[0]); widths=[15.0*cm/col_count]*col_count
            add_table(rows,widths)
        continue
    if line.startswith('- '):
        story.append(Paragraph(line[2:].replace('**',''), styles['BulletZ'], bulletText='•')); i+=1; continue
    if re.match(r'^\d+\. ', line):
        story.append(Paragraph(re.sub(r'^\d+\. ', '', line).replace('**',''), styles['BulletZ'], bulletText='•')); i+=1; continue
    story.append(Paragraph(line.replace('**',''), styles['BodyJust'])); i+=1


def footer(canvas, doc):
    canvas.saveState()
    canvas.setFont('Helvetica',7)
    canvas.setFillColor(colors.HexColor('#66657A'))
    canvas.drawString(1.2*cm, 1.08*cm, 'SKUs: ZT_STIAS_PROSBC_SETUP_CLIENT_A | ZT_PROSBC_HW500_STIAS_HYBRID_API_A | cupom hybrid7800')
    canvas.drawCentredString(A4[0]/2, 0.75*cm, f'ZICTEC · Chamada Verificada · Página {doc.page}')
    canvas.restoreState()

pdfdoc=SimpleDocTemplate(str(PDF), pagesize=A4, rightMargin=1.6*cm, leftMargin=1.6*cm, topMargin=1.5*cm, bottomMargin=1.4*cm)
pdfdoc.build(story, onFirstPage=footer, onLaterPages=footer)
shutil.copy2(PDF, CACHE/PDF.name)

# Validation
assert PPTX.exists() and PPTX.stat().st_size > 100000
assert DOCX.exists() and DOCX.stat().st_size > 20000
assert PDF.exists() and PDF.stat().st_size > 10000
prs = Presentation(str(PPTX))
assert len(prs.slides) == 11, len(prs.slides)
texts=[]
for s in prs.slides:
    for sh in s.shapes:
        if hasattr(sh, 'text'):
            texts.append(sh.text)
ppt_text='\n'.join(texts)
ppt_norm = ppt_text.replace('\n', '').replace(' ', '')
for kw in ['Origem Verificada','API REST Tools','Starter AS-OOB SaaS','hybrid7800']:
    assert kw in ppt_text, kw
for kw in ['ZT_STIAS_PROSBC_SETUP_CLIENT_A', 'ZT_PROSBC_HW500_STIAS_HYBRID_API_A']:
    assert kw.replace(' ', '') in ppt_norm, kw
Document(str(DOCX))
with zipfile.ZipFile(PPTX) as z:
    assert '[Content_Types].xml' in z.namelist()
with zipfile.ZipFile(DOCX) as z:
    xml=z.read('word/document.xml').decode('utf-8', errors='ignore')
    for kw in ['Origem Verificada','ZT_STIAS_API_A','ZT_PROSBC_HW500_STIAS_HYBRID_API_A']:
        assert kw in xml, kw

# Approximate PNG previews from PPTX shapes for visual QA fallback when LibreOffice is unavailable.
# It draws slide backgrounds, pictures and text positions from the actual PPTX object model.
W,H=1600,900
SX=W/prs.slide_width; SY=H/prs.slide_height
try:
    font_big=ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf', 32)
    font_med=ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf', 22)
    font_small=ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf', 18)
    font_tiny=ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf', 13)
except Exception:
    font_big=font_med=font_small=font_tiny=None

from pptx.enum.shapes import MSO_SHAPE_TYPE

def rgb_tuple(rgb, default=(255,255,255)):
    if rgb is None: return default
    try: return (rgb[0],rgb[1],rgb[2])
    except Exception: return default

preview_paths=[]
for idx, slide in enumerate(prs.slides, start=1):
    im=Image.new('RGB',(W,H),(251,248,246)); d=ImageDraw.Draw(im)
    # simple bg from first full-bleed rectangle fill
    for sh in slide.shapes:
        try:
            if sh.left == 0 and sh.top == 0 and sh.width >= prs.slide_width*0.9 and sh.height >= prs.slide_height*0.9 and sh.fill.type:
                c = sh.fill.fore_color.rgb
                im.paste(rgb_tuple(c,(251,248,246)), [0,0,W,H])
                break
        except Exception: pass
    for sh in slide.shapes:
        x=int(sh.left*SX); y=int(sh.top*SY); w=int(sh.width*SX); h=int(sh.height*SY)
        if w<=2 or h<=2: continue
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob=sh.image.blob
                import io
                pic=Image.open(io.BytesIO(blob)).convert('RGB')
                pic.thumbnail((w,h))
                px=x+(w-pic.width)//2; py=y+(h-pic.height)//2
                im.paste(pic,(px,py))
            except Exception: pass
        elif hasattr(sh, 'text') and sh.text.strip():
            text=sh.text.strip().replace('\x0b',' ')
            size = font_big if h>80 and w>500 and len(text)<95 else font_med if len(text)<45 else font_small if len(text)<160 else font_tiny
            # text color best effort
            color=(23,23,41)
            try:
                r=sh.text_frame.paragraphs[0].runs[0]
                color=rgb_tuple(r.font.color.rgb,color)
            except Exception: pass
            # wrap
            max_chars=max(10,int(w/(10 if size==font_small else 14)))
            words=text.split(); lines=[]; cur=''
            for word in words:
                if len(cur)+len(word)+1>max_chars:
                    lines.append(cur); cur=word
                else: cur=(cur+' '+word).strip()
            if cur: lines.append(cur)
            yy=y+6
            for ln in lines[:max(1,h//22)]:
                d.text((x+6,yy), ln, fill=color, font=size); yy+=24 if size!=font_tiny else 17
    out=PREVIEW_DIR / f'slide-{idx:02d}.png'
    im.save(out)
    preview_paths.append(out)
# contact sheet
thumbs=[Image.open(p).resize((320,180)) for p in preview_paths]
sheet=Image.new('RGB',(320*3,180*4),(240,240,240)); dr=ImageDraw.Draw(sheet)
for i,t in enumerate(thumbs):
    x=(i%3)*320; y=(i//3)*180
    sheet.paste(t,(x,y)); dr.text((x+6,y+6),str(i+1),fill=(255,90,31),font=font_med)
contact=PREVIEW_DIR/'pptx-contact-sheet.png'
sheet.save(contact)
shutil.copy2(contact, CACHE/contact.name)

print('QA_OK')
print('slides', len(prs.slides))
print('pptx_bytes', PPTX.stat().st_size)
print('docx_bytes', DOCX.stat().st_size)
print('pdf_bytes', PDF.stat().st_size)
print('preview', contact)
print('cache_pdf', CACHE/PDF.name)
print('keywords_ok')
