from pathlib import Path
from datetime import datetime
import shutil
import zipfile
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from docx import Document
from docx.shared import Cm, Pt as DPt, RGBColor as DRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / 'assets'
OUT = ROOT / 'artifacts' / 'chamada_verificada_collateral'
CACHE = Path('/root/.hermes/document_cache')
OUT.mkdir(parents=True, exist_ok=True)
CACHE.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(27, 22, 72)
BLUE = RGBColor(36, 33, 95)
ORANGE = RGBColor(255, 90, 31)
ORANGE2 = RGBColor(255, 122, 61)
SOFT = RGBColor(255, 241, 234)
BG = RGBColor(251, 248, 246)
INK = RGBColor(23, 23, 41)
MUTED = RGBColor(102, 101, 122)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(17, 163, 106)
LINE = RGBColor(235, 228, 223)

TODAY = datetime.now().strftime('%d/%m/%Y')

# ---------- Shared content ----------
slides = [
    {
        'kicker': 'ZICTEC · STIR/SHAKEN · Origem Verificada',
        'title': 'Autenticação de chamadas para reduzir risco e preparar Origem Verificada',
        'body': 'Pacotes comerciais e técnicos para sair do conceito para implantação real: API, ProSBC, ISBC Cloud, homologação e go-live assistido.',
        'image': 'hero.png',
    },
    {
        'title': 'O problema: confiança da chamada virou requisito operacional',
        'cards': [
            ('Risco na origem', 'Uso indevido de números, contestação e bloqueios exigem rastreabilidade.'),
            ('Regulação e ecossistema', 'Autenticação depende de aderência técnica, credenciais e homologação.'),
            ('Compra incompleta', 'Setup, licença, API, ISBC e hardware precisam ser separados para evitar gaps.'),
        ]
    },
    {
        'title': 'Conceitos essenciais',
        'flow': [
            ('STIR/SHAKEN', 'Protocolo/tecnologia de assinatura e validação da origem.'),
            ('Autenticação', 'Envio de chamadas em interconexões autenticadas conforme regras operacionais.'),
            ('Origem Verificada', 'Camada de identificação/branded call com marca, logo e campanha.'),
        ],
        'note': 'Identificação depende de Autenticação. Primeiro a chamada precisa ser confiável; depois pode receber camada de marca/campanha.'
    },
    {
        'title': 'Ecossistema ZICTEC: quatro formas de consumir',
        'products': [
            ('SaaS + SIP', 'Starter AS-OOB SaaS', 'ISBC Cloud para suporte SIP consumindo API ZICTEC.', 'stir-shaken-starter-as-oob-saas-anual-main.png'),
            ('API sem SIP', 'API REST Tools', 'Licença anual para SBCs/softswitches compatíveis.', 'stir-shaken-api-rest-tools-zictec.png'),
            ('Client ProSBC', 'Client local por SBC', 'Script/cliente proprietário instalado no ProSBC.', 'stir-shaken-client-prosbc-zictec.png'),
            ('Bundle híbrido', 'API + Client ProSBC', 'Combina API REST SaaS + cliente ProSBC.', 'stir-shaken-bundle-hibrido-api-client-prosbc.png'),
        ]
    },
    {
        'title': 'Oferta em caminhos simples',
        'cards': [
            ('A · Setup pontual', 'Implantação, UAT/homologação e go-live. Serviço, não licença.'),
            ('B · SaaS + SIP', 'Starter AS-OOB SaaS com ISBC Cloud, setup promocionalmente isento.'),
            ('C · API direta', 'Tools/API REST para integração direta, sem SIP/ISBC.'),
            ('D · Client ProSBC', 'Licença anual por SBC; no primeiro SBC normalmente combina com setup.'),
            ('Pacote 1º SBC', 'SETUP + Client ProSBC 1 ano com desconto promocional.'),
            ('Modernize', 'Servidor + 500 sessões + setup + client quando precisa base funcional.'),
        ]
    },
    {
        'title': 'Matriz de recomendação comercial',
        'matrix': [
            ('Precisa SIP/ISBC', 'SAAS_STIAS_OOB_A', 'Starter AS-OOB SaaS'),
            ('Consumo REST direto', 'ZT_STIAS_API_A', 'API REST Tools'),
            ('ProSBC adicional', 'ZT_STIAS_PROSBC_CLIENT_A', 'Client ProSBC por SBC'),
            ('1º ProSBC sem servidor', 'ZT_STIAS_PROSBC_SETUP_CLIENT_A', 'Setup + Client'),
            ('Servidor funcional', 'ZT_PROSBC_HW500_STIAS_SETUP_CLIENT_A', 'Server + 500 + Setup + Client'),
            ('Híbrido completo', 'ZT_PROSBC_HW500_STIAS_HYBRID_API_A', 'Server + Client + API Tools'),
        ]
    },
    {
        'title': 'Jornada ideal do cliente',
        'steps': ['Educar em 2 minutos', 'Escolher objetivo', 'Responder checklist técnico', 'Receber pacote recomendado', 'Checkout ou lead qualificado'],
        'image': 'technical.png'
    },
    {
        'title': 'Checklist de contratação / kickoff',
        'cards': [
            ('Cliente e regulatório', 'Razão social, CNPJ, responsável técnico, status ABR e credenciais.'),
            ('Topologia e tráfego', 'SBC/vendor, rotas, volumetria, IN-BAND/OUT-OF-BAND e hardware.'),
            ('Operação e aceite', 'Relatórios, suporte pós-go-live, premissas fora de escopo e critério de aceite.'),
        ]
    },
    {
        'title': 'Pacotes e valores de referência',
        'pricing': [
            ('SAAS_STIAS_OOB_A', 'Starter AS-OOB SaaS', 'R$ 39.000,00'),
            ('ZT_STIAS_API_A', 'API REST Tools', 'R$ 31.200,00'),
            ('ZT_STIAS_PROSBC_CLIENT_A', 'Client ProSBC anual', 'R$ 12.000,00'),
            ('ZT_STIAS_PROSBC_SETUP_CLIENT_A', 'SETUP + Client 1º ano', 'R$ 18.400,04'),
            ('ZT_PROSBC_HW500_STIAS_SETUP_CLIENT_A', 'Server + 500 + SETUP + Client', 'R$ 33.864,44'),
            ('ZT_PROSBC_HW500_STIAS_HYBRID_API_A', 'Server + Client + API Tools', 'R$ 49.464,44'),
        ]
    },
    {
        'title': 'Premissas e limites de escopo',
        'cards': [
            ('Compra inicia validação', 'Produção depende de ambiente, rotas, credenciais, homologação e aceite.'),
            ('Setup não é licença', 'Para ProSBC, setup precisa ser combinado com Client ou pacote fechado.'),
            ('API ≠ ISBC Cloud', 'REST direto, SIP via ISBC e Client ProSBC são canais diferentes.'),
            ('Servidor físico é variação', 'Pacotes com hardware devem deixar claro o que renova no segundo ano.'),
        ]
    },
    {
        'title': 'Próximo passo comercial',
        'cards': [
            ('1', 'Usar o configurador para pré-qualificar cenário e SKU.'),
            ('2', 'Copiar briefing para pré-venda/técnico antes do checkout.'),
            ('3', 'Direcionar compra direta, validação obrigatória ou proposta formal.'),
        ],
        'cta': 'https://jarvis-zictec.github.io/chamada-verificada-jornada/'
    },
]

whitepaper_md = f"""# Chamada Verificada ZICTEC: Datasheet técnico-comercial

**Fonte:** página Chamada Verificada ZICTEC validada em GitHub Pages  
**Data:** {TODAY}  
**Uso recomendado:** material comercial, pré-venda, reunião executiva e alinhamento de escopo técnico.

## Resumo executivo

A jornada de Chamada Verificada da ZICTEC organiza a oferta de Autenticação de Chamadas, STIR/SHAKEN e preparação para Origem Verificada em pacotes claros. O objetivo é reduzir risco operacional, evitar compra incompleta e guiar o cliente entre consumo por SIP/ISBC, API REST, Client ProSBC, bundle híbrido e pacotes com servidor físico.

## Conceitos centrais

- **STIR/SHAKEN:** tecnologia/protocolo usado para assinatura e validação da origem da chamada.
- **Autenticação de chamadas:** envio das chamadas nas interconexões autenticadas conforme regras do ecossistema operacional e regulatório.
- **Origem Verificada:** camada brasileira de identificação/branded call, com marca, logo e campanha. Depende de Autenticação.
- **ABR Telecom:** referência institucional e operacional a ser usada na terminologia comercial.

## Benefícios para o cliente

1. **Menos risco na origem da chamada:** validação técnica, trilha de auditoria e critérios para reduzir uso indevido de números, falhas de autenticação e contestação.
2. **Pacote certo para a topologia atual:** separa ISBC Cloud, API REST, Client ProSBC, setup, licença, hardware e suporte.
3. **Base para Origem Verificada:** primeiro autenticação; depois identificação, marca e campanhas com dependências explícitas.

## Formas de consumir a solução

| Modalidade | Quando usar | O que entrega | Atenções |
|---|---|---|---|
| Starter AS-OOB SaaS | Cliente precisa de SIP/ISBC Cloud | ISBC Cloud compartilhado consumindo API ZICTEC | Setup promocionalmente isento; validar limites Light |
| API REST ZICTEC Tools | Cliente compatível quer consumo REST direto | API, relatórios, auditoria, CDR, volumetria e falhas | Não inclui ISBC/SIP nem Client ProSBC |
| Client ProSBC | Integração precisa rodar no ProSBC | Script/cliente proprietário licenciado por SBC | No primeiro SBC normalmente combina com setup |
| Bundle híbrido | Cliente ProSBC quer API + client | API REST SaaS + Client ProSBC | Bundle anual sem servidor não inclui setup |
| Pacotes com servidor | Cliente precisa base ProSBC funcional | Hardware, 500 sessões, setup/client e, no híbrido, API Tools | Suporte operacional e expansões ficam fora salvo contratação |

## Matriz de pacotes / SKUs

| Cenário | SKU / pacote | Valor referência | Leitura comercial |
|---|---:|---:|---|
| SIP via ISBC compartilhado | SAAS_STIAS_OOB_A | R$ 39.000,00 | Starter AS-OOB SaaS, plano Light, setup isento na promoção |
| API REST direto | ZT_STIAS_API_A | R$ 31.200,00 | Tools/relatórios/API para cliente compatível, sem SIP |
| Client ProSBC adicional | ZT_STIAS_PROSBC_CLIENT_A | R$ 12.000,00 | Licença anual por SBC, não inclui setup |
| Primeiro ProSBC sem servidor | ZT_STIAS_PROSBC_SETUP_CLIENT_A | R$ 18.400,04 | SETUP + Client 1 ano, 30% off sobre setup+client |
| Servidor + ProSBC funcional | ZT_PROSBC_HW500_STIAS_SETUP_CLIENT_A | R$ 33.864,44 | Hardware + 500 sessões + SETUP + Client |
| Híbrido completo 1º ano | ZT_PROSBC_HW500_STIAS_HYBRID_API_A | R$ 49.464,44 | Servidor + 500 sessões + Client + API Tools; cupom hybrid7800 |
| Setup avulso | ZTSHAKENDEPLOY | R$ 14.285,78 | Serviço de implantação/homologação, não é licença |

**Observação promocional:** o pacote ZT_PROSBC_HW500_STIAS_HYBRID_API_A usa o cupom promocional **hybrid7800** no shop, conforme jornada comercial publicada.

## Jornada de contratação

1. **Educar em 2 minutos:** diferenciar STIR/SHAKEN, Autenticação e Origem Verificada.
2. **Escolher objetivo:** cumprir regulação, operar com suporte, manter ProSBC atual, modernizar topologia ou ativar campanhas.
3. **Responder checklist técnico:** SBC/vendor, rotas, volumetria, status ABR, relatórios, campanhas e suporte.
4. **Receber pacote recomendado:** API REST, ISBC/SIP, Client ProSBC, bundle híbrido ou pacote com servidor.
5. **Checkout ou lead qualificado:** compra inicia validação; cenários complexos viram proposta formal.

## Checklist de kickoff

### Cliente e regulatório
- Razão social, CNPJ e responsável técnico.
- Status ABR Telecom: não iniciado, UAT/homologação ou produção.
- Credenciais/certificados disponíveis.
- Objetivo: Autenticação, Identificação/Origem Verificada ou ambos.

### Topologia e tráfego
- SBC/vendor, versão e capacidade.
- Rotas/interconexões e volumetria.
- IN-BAND, OUT-OF-BAND ou indefinido.
- Necessidade de ProSBC, licença, servidor físico ou expansão.

### Operação e aceite
- Relatórios/auditoria esperados.
- Suporte pós-go-live e SLA desejado.
- Premissas fora de escopo.
- Critério de homologação e aceite.

## Premissas e exclusões comerciais

- O ecommerce deve reduzir fricção, não prometer implantação automática universal.
- A compra inicia contratação/validação; produção depende de ambiente, rotas, credenciais, homologação e aceite.
- API REST Tools não é ISBC Cloud; Client ProSBC não é licença da API SaaS.
- SETUP é serviço, não licença. Para ProSBC, combinar com Client ou usar pacote fechado.
- Pacotes com servidor físico são variações vinculadas aos itens sem servidor; o que renova no segundo ano deve ficar explícito.
- Suporte operacional, banco de horas, transcoding, expansões e SLA dedicado ficam fora salvo contratação separada.

## Próxima ação recomendada

Usar o configurador comercial para gerar briefing de pré-venda e selecionar a rota adequada: compra direta, validação obrigatória ou proposta formal. Para reuniões comerciais, anexar este datasheet e a apresentação executiva como materiais de apoio.
"""

# ---------- PPTX helpers ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]

def add_bg(slide, color=BG):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.z_order = 0


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, font='Aptos', valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(0.02)
    box.text_frame.margin_right = Inches(0.02)
    box.text_frame.margin_top = Inches(0.02)
    box.text_frame.margin_bottom = Inches(0.02)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return box


def add_round_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(radius, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    return shp


def add_card(slide, x, y, w, h, title, body, accent=False):
    shp = add_round_rect(slide, x, y, w, h, fill=NAVY if accent else WHITE, line=NAVY if accent else LINE)
    add_text(slide, title, x+0.25, y+0.22, w-0.5, 0.45, 17, WHITE if accent else NAVY, True)
    add_text(slide, body, x+0.25, y+0.82, w-0.5, h-1.0, 12.5, WHITE if accent else MUTED)
    return shp


def add_header(slide, title, kicker=None, dark=False):
    if kicker:
        add_text(slide, kicker.upper(), 0.65, 0.35, 6, 0.25, 9.5, ORANGE if not dark else RGBColor(255,178,153), True)
    add_text(slide, title, 0.62, 0.62, 8.4, 0.85, 29, WHITE if dark else NAVY, True)
    try:
        slide.shapes.add_picture(str(ASSETS/'logo_orange.png' if dark else ASSETS/'logo_blue.png'), Inches(11.4), Inches(0.35), width=Inches(1.2))
    except Exception:
        add_text(slide, 'ZICTEC', 11.4, 0.35, 1.2, 0.3, 16, ORANGE if dark else NAVY, True)


def add_footer(slide, dark=False):
    add_text(slide, 'ZICTEC · Chamada Verificada · material técnico-comercial', 0.65, 7.1, 8, 0.22, 8.5, RGBColor(180,180,198) if dark else MUTED)

# Slide generation
for i, spec in enumerate(slides):
    slide = prs.slides.add_slide(blank)
    dark = i in (0, len(slides)-1)
    add_bg(slide, NAVY if dark else BG)
    add_header(slide, spec['title'], spec.get('kicker'), dark=dark)
    add_footer(slide, dark=dark)

    if i == 0:
        add_text(slide, spec['body'], 0.68, 1.85, 5.9, 1.0, 18, WHITE, False)
        for idx, (label, txt) in enumerate([('Start','cumprir requisito regulatório'),('Operate','SaaS/ISBC e relatórios'),('Hybrid','ProSBC + API'),('Modernize','servidor e expansão')]):
            x = 0.68 + (idx%2)*2.65; y = 3.25 + (idx//2)*1.0
            add_round_rect(slide, x, y, 2.35, 0.72, fill=BLUE, line=BLUE)
            add_text(slide, label, x+0.15, y+0.12, 0.9, 0.25, 14, ORANGE, True)
            add_text(slide, txt, x+0.15, y+0.38, 1.95, 0.25, 9.5, WHITE)
        img = ASSETS / spec['image']
        slide.shapes.add_picture(str(img), Inches(7.05), Inches(1.25), width=Inches(5.4))
        add_round_rect(slide, 7.35, 5.7, 4.7, 0.82, fill=RGBColor(44,38,107), line=RGBColor(44,38,107))
        add_text(slide, 'Princípio central: Origem Verificada depende de Autenticação.', 7.62, 5.93, 4.1, 0.3, 13.5, WHITE, True)
    elif 'cards' in spec and i in (1,7,9,10):
        cards = spec['cards']
        if len(cards) == 3:
            for j, (t,b) in enumerate(cards):
                add_card(slide, 0.75 + j*4.1, 2.05, 3.65, 2.65, t, b, accent=(j==1 and i==1))
        else:
            for j, (t,b) in enumerate(cards):
                add_card(slide, 0.75 + (j%2)*6.0, 1.75 + (j//2)*1.75, 5.4, 1.25, t, b, accent=(j==0))
    elif 'flow' in spec:
        for j, (t,b) in enumerate(spec['flow']):
            x = 0.8 + j*4.15
            add_round_rect(slide, x, 2.1, 3.35, 2.2, fill=WHITE if j != 1 else NAVY, line=ORANGE if j == 1 else LINE)
            add_text(slide, str(j+1), x+0.22, 2.32, 0.45, 0.45, 17, ORANGE, True, align=PP_ALIGN.CENTER)
            add_text(slide, t, x+0.78, 2.25, 2.3, 0.35, 18, NAVY if j != 1 else WHITE, True)
            add_text(slide, b, x+0.28, 2.95, 2.8, 0.9, 12, MUTED if j != 1 else WHITE)
            if j < 2:
                add_text(slide, '→', x+3.55, 2.85, 0.45, 0.35, 30, ORANGE, True, align=PP_ALIGN.CENTER)
        add_round_rect(slide, 1.2, 5.15, 10.9, 0.85, fill=SOFT, line=SOFT)
        add_text(slide, spec['note'], 1.45, 5.42, 10.4, 0.25, 14, INK, True)
    elif 'products' in spec:
        for j,(badge,t,b,img) in enumerate(spec['products']):
            x = 0.62 + j*3.16
            add_round_rect(slide, x, 1.75, 2.78, 4.6, fill=WHITE, line=ORANGE if j==0 else LINE)
            slide.shapes.add_picture(str(ASSETS/img), Inches(x+0.25), Inches(2.05), width=Inches(2.25), height=Inches(1.55))
            add_text(slide, badge.upper(), x+0.27, 3.83, 2.2, 0.22, 8.5, ORANGE, True)
            add_text(slide, t, x+0.27, 4.15, 2.25, 0.55, 15, NAVY, True)
            add_text(slide, b, x+0.27, 4.85, 2.25, 0.9, 10.8, MUTED)
    elif i == 4:
        for j,(t,b) in enumerate(spec['cards']):
            add_card(slide, 0.7+(j%3)*4.2, 1.62+(j//3)*2.05, 3.65, 1.55, t, b, accent=j in (1,4))
    elif 'matrix' in spec:
        rows = spec['matrix']
        table = slide.shapes.add_table(len(rows)+1,3,Inches(0.75),Inches(1.65),Inches(11.85),Inches(4.95)).table
        widths=[3.6,3.7,4.55]
        for c,w in enumerate(widths): table.columns[c].width=Inches(w)
        headers=['Sinal do cliente','SKU / pacote','Oferta sugerida']
        for c,h in enumerate(headers):
            cell=table.cell(0,c); cell.text=h; cell.fill.solid(); cell.fill.fore_color.rgb=NAVY
            for p in cell.text_frame.paragraphs:
                for r in p.runs: r.font.color.rgb=WHITE; r.font.bold=True; r.font.size=Pt(11)
        for r,row in enumerate(rows,1):
            for c,val in enumerate(row):
                cell=table.cell(r,c); cell.text=val
                cell.fill.solid(); cell.fill.fore_color.rgb = SOFT if r%2==0 else WHITE
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size=Pt(9.2); run.font.color.rgb=INK if c!=1 else ORANGE; run.font.bold=(c==1)
    elif 'steps' in spec:
        for j,step in enumerate(spec['steps']):
            y=1.55+j*0.92
            add_round_rect(slide, 0.8, y, 0.55, 0.55, fill=ORANGE, line=ORANGE)
            add_text(slide, str(j+1), 0.98, y+0.14, 0.2, 0.2, 12, WHITE, True, align=PP_ALIGN.CENTER)
            add_text(slide, step, 1.55, y+0.12, 4.4, 0.3, 16, NAVY, True)
        slide.shapes.add_picture(str(ASSETS/spec['image']), Inches(7.0), Inches(1.55), width=Inches(5.1))
    elif 'pricing' in spec:
        rows=spec['pricing']
        for j,(sku,name,price) in enumerate(rows):
            x=0.75+(j%2)*6.05; y=1.55+(j//2)*1.55
            add_round_rect(slide, x, y, 5.55, 1.18, fill=WHITE, line=ORANGE if j in (3,5) else LINE)
            add_text(slide, sku, x+0.25, y+0.18, 2.95, 0.25, 10.5, ORANGE, True)
            add_text(slide, name, x+0.25, y+0.53, 3.15, 0.3, 12.5, NAVY, True)
            add_text(slide, price, x+3.55, y+0.38, 1.65, 0.35, 15, NAVY, True, align=PP_ALIGN.RIGHT)
        add_text(slide, 'Nota: pacote híbrido 1º ano usa cupom promocional hybrid7800 conforme página do shop.', 0.78, 6.55, 11.6, 0.28, 10, ORANGE, True, align=PP_ALIGN.CENTER)
    elif i == 10:
        for j,(n,b) in enumerate(spec['cards']):
            x=1.1+j*3.75
            add_round_rect(slide, x, 2.1, 2.95, 2.25, fill=WHITE, line=ORANGE)
            add_text(slide, n, x+1.1, 2.28, 0.7, 0.5, 24, ORANGE, True, align=PP_ALIGN.CENTER)
            add_text(slide, b, x+0.28, 3.02, 2.4, 0.8, 13, INK, True, align=PP_ALIGN.CENTER)
        add_round_rect(slide, 2.5, 5.35, 8.3, 0.65, fill=ORANGE, line=ORANGE)
        add_text(slide, spec['cta'], 2.8, 5.55, 7.6, 0.2, 11.5, WHITE, True, align=PP_ALIGN.CENTER)

pptx_path = OUT / 'zictec-chamada-verificada-apresentacao.pptx'
prs.save(pptx_path)

# ---------- DOCX ----------
md_path = OUT / 'zictec-chamada-verificada-datasheet-whitepaper.md'
md_path.write_text(whitepaper_md, encoding='utf-8')

def shade_cell(cell, fill):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:fill'), fill)
    tcPr.append(shd)


def set_cell_text(cell, text, bold=False, color='171729', size=9):
    cell.text = ''
    p = cell.paragraphs[0]
    run = p.add_run(text)
    run.bold = bold
    run.font.size = DPt(size)
    run.font.color.rgb = DRGBColor.from_string(color)


doc = Document()
sec = doc.sections[0]
sec.top_margin = Cm(1.8); sec.bottom_margin = Cm(1.8); sec.left_margin = Cm(2.0); sec.right_margin = Cm(2.0)
styles = doc.styles
styles['Normal'].font.name = 'Aptos'
styles['Normal'].font.size = DPt(10.5)
styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), 'Aptos')
for style_name, size, color in [('Heading 1',18,'1B1648'),('Heading 2',14,'1B1648'),('Heading 3',12,'1B1648')]:
    st=styles[style_name]; st.font.name='Aptos'; st.font.size=DPt(size); st.font.color.rgb=DRGBColor.from_string(color); st.font.bold=True

header = sec.header.paragraphs[0]
header.text = 'ZICTEC | Chamada Verificada'
header.alignment = WD_ALIGN_PARAGRAPH.RIGHT
header.runs[0].font.color.rgb = DRGBColor.from_string('1B1648')
header.runs[0].font.bold = True
footer = sec.footer.paragraphs[0]
footer.text = 'Documento técnico-comercial • material de apoio à pré-venda'
footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
footer.runs[0].font.size = DPt(8)
footer.runs[0].font.color.rgb = DRGBColor.from_string('66657A')

# cover
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = p.add_run('ZICTEC')
r.font.size = DPt(22); r.font.bold=True; r.font.color.rgb = DRGBColor.from_string('FF5A1F')
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = p.add_run('DATASHEET / WHITEPAPER TÉCNICO-COMERCIAL')
r.font.size = DPt(10); r.font.bold=True; r.font.color.rgb = DRGBColor.from_string('66657A')
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = p.add_run('Chamada Verificada, STIR/SHAKEN e Origem Verificada')
r.font.size = DPt(24); r.font.bold=True; r.font.color.rgb = DRGBColor.from_string('1B1648')
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = p.add_run('Guia de ofertas, arquitetura comercial, matriz de decisão e premissas para implantação assistida.')
r.font.size = DPt(12); r.font.color.rgb = DRGBColor.from_string('444157')

meta = doc.add_table(rows=3, cols=2)
meta.alignment = WD_TABLE_ALIGNMENT.CENTER
for row in meta.rows:
    for cell in row.cells:
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
pairs=[('Data',TODAY),('Preparado por','Jarvis / ZICTEC'),('Classificação','Material técnico-comercial')]
for i,(k,v) in enumerate(pairs):
    set_cell_text(meta.cell(i,0),k,True,'FFFFFF',9); shade_cell(meta.cell(i,0),'1B1648')
    set_cell_text(meta.cell(i,1),v,False,'171729',9); shade_cell(meta.cell(i,1),'FFF1EA')
doc.add_page_break()

# executive callout
h = doc.add_heading('Resumo executivo', level=1)
summary_table = doc.add_table(rows=1, cols=1)
summary_cell = summary_table.cell(0,0)
shade_cell(summary_cell, 'FFF1EA')
set_cell_text(summary_cell, 'A ZICTEC organiza Chamada Verificada em uma jornada guiada: educa o cliente, separa Autenticação de Origem Verificada, diferencia API REST, SIP/ISBC, Client ProSBC, bundle híbrido e pacotes com servidor físico, e reduz risco de compra incompleta antes do checkout.', True, '1B1648', 10)

# Simple markdown renderer
lines = whitepaper_md.splitlines()
skip_h1 = True
for line in lines:
    line=line.rstrip()
    if not line:
        continue
    if line.startswith('# '):
        if skip_h1:
            skip_h1=False
            continue
        doc.add_heading(line[2:], level=1)
    elif line.startswith('## '):
        if line[3:] == 'Resumo executivo':
            continue
        doc.add_heading(line[3:], level=1)
    elif line.startswith('### '):
        doc.add_heading(line[4:], level=2)
    elif line.startswith('|'):
        continue
    elif line.startswith('- '):
        p=doc.add_paragraph(style='List Bullet')
        p.add_run(line[2:].replace('**','').replace('`',''))
    elif re.match(r'^\d+\. ', line):
        p=doc.add_paragraph(style='List Number')
        p.add_run(re.sub(r'^\d+\. ', '', line).replace('**','').replace('`',''))
    elif line.startswith('**') and line.endswith('**'):
        p=doc.add_paragraph()
        r=p.add_run(line.strip('*'))
        r.bold=True
    else:
        # skip duplicated source metadata on cover? keep if not too redundant
        p=doc.add_paragraph()
        txt=line.replace('**','').replace('`','')
        p.add_run(txt)
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

# Add explicit tables (Markdown tables are skipped above)
doc.add_page_break()
doc.add_heading('Tabela executiva de modalidades', level=1)
rows = [
    ('Starter AS-OOB SaaS','SIP/ISBC Cloud','Base SaaS operacional com relatórios/auditoria','Validar limites Light'),
    ('API REST Tools','REST direto','API, auditoria, CDR, volumetria e falhas','Não inclui SIP/ISBC'),
    ('Client ProSBC','ProSBC local','Script/cliente por SBC','Não inclui setup'),
    ('Bundle híbrido','API + Client','API SaaS + Client ProSBC','Sem servidor e sem setup'),
    ('Pacotes com servidor','Modernização','Hardware + 500 sessões + STIR/SHAKEN','Suporte fora salvo contratação'),
]
t=doc.add_table(rows=len(rows)+1, cols=4)
t.style='Table Grid'
headers=('Modalidade','Canal','Entrega','Atenção')
for c,h in enumerate(headers): set_cell_text(t.cell(0,c),h,True,'FFFFFF',8); shade_cell(t.cell(0,c),'1B1648')
for r,row in enumerate(rows,1):
    for c,val in enumerate(row): set_cell_text(t.cell(r,c),val, c==0, '171729' if c!=0 else 'FF5A1F',8)

doc.add_heading('Matriz de SKUs e valores de referência', level=1)
rows = [
    ('SAAS_STIAS_OOB_A','Starter AS-OOB SaaS','R$ 39.000,00'),
    ('ZT_STIAS_API_A','API REST Tools','R$ 31.200,00'),
    ('ZT_STIAS_PROSBC_CLIENT_A','Client ProSBC anual','R$ 12.000,00'),
    ('ZT_STIAS_PROSBC_SETUP_CLIENT_A','SETUP + Client 1º ano','R$ 18.400,04'),
    ('ZT_PROSBC_HW500_STIAS_SETUP_CLIENT_A','Server + 500 + SETUP + Client','R$ 33.864,44'),
    ('ZT_PROSBC_HW500_STIAS_HYBRID_API_A','Server + Client + API Tools','R$ 49.464,44'),
    ('ZTSHAKENDEPLOY','Setup avulso','R$ 14.285,78'),
]
t=doc.add_table(rows=len(rows)+1, cols=3)
t.style='Table Grid'
for c,h in enumerate(('SKU','Oferta','Valor')): set_cell_text(t.cell(0,c),h,True,'FFFFFF',8); shade_cell(t.cell(0,c),'1B1648')
for r,row in enumerate(rows,1):
    for c,val in enumerate(row):
        set_cell_text(t.cell(r,c),val, c==0, 'FF5A1F' if c==0 else '171729',8)

doc.add_paragraph()
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
rr = p.add_run('Próxima ação: usar o configurador comercial e anexar este material junto à apresentação executiva em reuniões de pré-venda.')
rr.bold=True; rr.font.color.rgb = DRGBColor.from_string('1B1648')

docx_path = OUT / 'zictec-chamada-verificada-datasheet-whitepaper.docx'
doc.save(docx_path)
# Reopen validation
Document(str(docx_path))

# Copy to document cache
for p in [pptx_path, docx_path, md_path]:
    shutil.copy2(p, CACHE / p.name)

print('ARTIFACTS')
for p in [pptx_path, docx_path, md_path, CACHE/pptx_path.name, CACHE/docx_path.name, CACHE/md_path.name]:
    print(p, p.stat().st_size)
